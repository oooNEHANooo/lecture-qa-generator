"""
PowerPoint content extraction service
"""

import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

logger = logging.getLogger(__name__)


class SlideContent:
    """スライドの内容を表すクラス"""
    
    def __init__(self, slide_number: int, title: str = "", content: str = "", 
                 bullet_points: List[str] = None, images: List[str] = None):
        self.slide_number = slide_number
        self.title = title
        self.content = content
        self.bullet_points = bullet_points or []
        self.images = images or []
        
    def to_dict(self) -> Dict[str, Any]:
        return {
            "slide_number": self.slide_number,
            "title": self.title,
            "content": self.content,
            "bullet_points": self.bullet_points,
            "images": self.images,
            "full_text": self.get_full_text()
        }
    
    def get_full_text(self) -> str:
        """スライドの全テキストを結合して返す"""
        parts = []
        if self.title:
            parts.append(f"タイトル: {self.title}")
        if self.content:
            parts.append(f"内容: {self.content}")
        if self.bullet_points:
            parts.append("要点:")
            for point in self.bullet_points:
                parts.append(f"• {point}")
        return "\n".join(parts)


class PPTXExtractor:
    """PowerPointファイルからコンテンツを抽出するサービス"""
    
    def __init__(self):
        self.presentation: Optional[Presentation] = None
        self.slides_content: List[SlideContent] = []
    
    def extract_from_file(self, file_path: Path) -> List[SlideContent]:
        """PowerPointファイルからコンテンツを抽出"""
        try:
            logger.info(f"PowerPointファイルを読み込み中: {file_path}")
            self.presentation = Presentation(file_path)
            self.slides_content = []
            
            for slide_number, slide in enumerate(self.presentation.slides, 1):
                slide_content = self._extract_slide_content(slide, slide_number)
                self.slides_content.append(slide_content)
                logger.debug(f"スライド {slide_number} を処理完了")
            
            logger.info(f"合計 {len(self.slides_content)} 枚のスライドを処理完了")
            return self.slides_content
            
        except Exception as e:
            logger.error(f"PowerPointファイルの読み込みに失敗: {e}")
            raise
    
    def _extract_slide_content(self, slide: Slide, slide_number: int) -> SlideContent:
        """個別のスライドからコンテンツを抽出"""
        title = ""
        content_parts = []
        bullet_points = []
        images = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = self._extract_text_from_shape(shape)
                if text:
                    # タイトルの判定（最初のテキストボックスまたはタイトル形状）
                    if not title and (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or 
                                     self._is_title_shape(shape)):
                        title = text
                    else:
                        # 箇条書きの判定
                        if self._is_bullet_point(text):
                            bullet_points.extend(self._extract_bullet_points(text))
                        else:
                            content_parts.append(text)
            
            # 画像の処理
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append(f"画像_{len(images) + 1}")
        
        content = "\n".join(content_parts)
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            content=content,
            bullet_points=bullet_points,
            images=images
        )
    
    def _extract_text_from_shape(self, shape: BaseShape) -> str:
        """図形からテキストを抽出"""
        if not shape.has_text_frame:
            return ""
        
        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = ""
            for run in paragraph.runs:
                paragraph_text += run.text
            if paragraph_text.strip():
                text_parts.append(paragraph_text.strip())
        
        return "\n".join(text_parts)
    
    def _is_title_shape(self, shape: BaseShape) -> bool:
        """図形がタイトルかどうかを判定"""
        if not shape.has_text_frame:
            return False
        
        # 位置とサイズでタイトルを判定（上部に配置され、大きなフォント）
        # 簡単な判定ロジック
        top = shape.top
        height = shape.height
        
        # スライドの上部30%以内にある場合はタイトルと判定
        return top < 914400 * 0.3  # 914400は1インチのEMU値
    
    def _is_bullet_point(self, text: str) -> bool:
        """テキストが箇条書きかどうかを判定"""
        bullet_indicators = ["•", "・", "-", "◦", "▪", "▫"]
        lines = text.split("\n")
        
        # 複数行があり、いずれかが箇条書き記号で始まる場合
        if len(lines) > 1:
            for line in lines:
                line = line.strip()
                if line and any(line.startswith(indicator) for indicator in bullet_indicators):
                    return True
        
        return False
    
    def _extract_bullet_points(self, text: str) -> List[str]:
        """箇条書きポイントを抽出"""
        bullet_indicators = ["•", "・", "-", "◦", "▪", "▫"]
        points = []
        
        lines = text.split("\n")
        for line in lines:
            line = line.strip()
            if line:
                # 箇条書き記号を除去
                for indicator in bullet_indicators:
                    if line.startswith(indicator):
                        line = line[len(indicator):].strip()
                        break
                
                if line:
                    points.append(line)
        
        return points
    
    def get_lecture_summary(self) -> Dict[str, Any]:
        """講義全体のサマリーを生成"""
        if not self.slides_content:
            return {}
        
        total_slides = len(self.slides_content)
        titles = [slide.title for slide in self.slides_content if slide.title]
        total_bullet_points = sum(len(slide.bullet_points) for slide in self.slides_content)
        total_images = sum(len(slide.images) for slide in self.slides_content)
        
        # 全体のテキスト量
        total_text_length = sum(len(slide.get_full_text()) for slide in self.slides_content)
        
        return {
            "total_slides": total_slides,
            "titles": titles,
            "total_bullet_points": total_bullet_points,
            "total_images": total_images,
            "total_text_length": total_text_length,
            "average_text_per_slide": total_text_length / total_slides if total_slides > 0 else 0
        }
    
    def get_slides_for_qa_generation(self) -> List[Dict[str, Any]]:
        """QA生成用のスライドデータを返す"""
        return [slide.to_dict() for slide in self.slides_content]