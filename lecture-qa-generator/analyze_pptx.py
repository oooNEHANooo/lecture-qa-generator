import os
from pptx import Presentation
import json

def analyze_pptx(pptx_path):
    """PPTXファイルの内容を解析する"""
    prs = Presentation(pptx_path)
    
    # 基本情報
    info = {
        'file_name': os.path.basename(pptx_path),
        'slide_count': len(prs.slides),
        'slides': []
    }
    
    # 各スライドの内容を取得
    for i, slide in enumerate(prs.slides):
        slide_info = {
            'slide_number': i + 1,
            'title': '',
            'content': []
        }
        
        # タイトルを取得
        if slide.shapes.title:
            slide_info['title'] = slide.shapes.title.text
        
        # すべてのテキストを取得
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:  # タイトル以外のテキスト
                    slide_info['content'].append(shape.text.strip())
        
        info['slides'].append(slide_info)
    
    return info

# Day1のPPTXを解析
day1_path = "PJT02_発注書_講義内容の確認QAの作成システム/LLM2023/Matsuo_Lab_LLM_Day1_20231227.pptx"
day1_info = analyze_pptx(day1_path)

print(f"ファイル名: {day1_info['file_name']}")
print(f"総スライド数: {day1_info['slide_count']}")
print("\n" + "="*60 + "\n")

# 最初の15スライドの内容を表示
for slide in day1_info['slides'][:15]:
    print(f"スライド {slide['slide_number']}:")
    if slide['title']:
        print(f"  タイトル: {slide['title']}")
    if slide['content']:
        print("  内容:")
        for content in slide['content'][:3]:  # 各スライドの最初の3つのコンテンツ
            if len(content) > 100:
                print(f"    - {content[:100]}...")
            else:
                print(f"    - {content}")
    print()

# 他のDay2-7のファイルも簡単に確認
print("\n" + "="*60 + "\n")
print("その他の講義資料の概要:")
for day in range(2, 8):
    file_path = f"PJT02_発注書_講義内容の確認QAの作成システム/LLM2023/Matsuo_Lab_LLM_Day{day}_20231227.pptx"
    try:
        prs = Presentation(file_path)
        print(f"\nDay{day}: {len(prs.slides)}スライド")
        # 最初のスライドのタイトルを表示
        if prs.slides and prs.slides[0].shapes.title:
            print(f"  最初のスライド: {prs.slides[0].shapes.title.text}")
    except Exception as e:
        print(f"\nDay{day}: エラー - {e}")